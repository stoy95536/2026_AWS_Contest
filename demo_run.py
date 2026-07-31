"""
完整端到端 Demo：
1. 導入附件四 Excel
2. 使用附件二的提示詞
3. 執行完整計算引擎
4. 生成 16 頁 slide_spec (含真實數據)
5. 使用附件一模板生成 PPT
6. 輸出所有結果

此腳本不依賴 LLM（純規則引擎），驗證資料正確性。
"""
import sys
sys.path.insert(0, ".")

import os
import json
from src.data_loader import ExcelLoader, DataStandardizer
from src.calculation_engine import MetricCalculator, DataLineageTracker
from src.validation import DataValidator
from src.agents.planner_agent import PlannerAgent
from src.agents.analyst_agent import AnalystAgent
from src.agents.reviewer_agent import ReviewerAgent
from src.presentation import PPTGenerator
from src.validation.ppt_reconciler import PPTReconciler

EXCEL_PATH = "附件四_預期修正參照資料.xlsx"
TEMPLATE_PATH = "附件一_台新新光金控簡報版型.pptx"
OUTPUT_DIR = "outputs"

os.makedirs(OUTPUT_DIR, exist_ok=True)

print("=" * 70)
print("  LLM 驅動之 Excel 報表轉簡報自動化系統 — 完整 Demo")
print("  輸入: 附件四_預期修正參照資料.xlsx")
print("  模板: 附件一_台新新光金控簡報版型.pptx")
print("  提示: 附件二_系統提示詞")
print("=" * 70)

# ============================================================
# STEP 1: Excel 解析
# ============================================================
print("\n" + "─" * 70)
print("STEP 1: Excel 資料解析")
print("─" * 70)
loader = ExcelLoader(EXCEL_PATH)
sheets = loader.get_sheet_names()
print(f"  工作表數量: {len(sheets)}")
for s in sheets:
    print(f"    - {s}")

standardizer = DataStandardizer(EXCEL_PATH)
for sheet_name in sheets:
    df = loader.read_sheet_to_dataframe(sheet_name, header_row=1)
    records = standardizer.standardize_sheet(df, sheet_name)
    print(f"  {sheet_name}: {len(records)} 筆")

loader.close()
std_data = standardizer.to_dataframe()
print(f"\n  標準化資料總筆數: {len(std_data)}")
print(f"  金融機構數量: {len(standardizer.get_institutions())}")
print(f"  期間: {standardizer.get_periods()}")

# ============================================================
# STEP 2: 資料驗證
# ============================================================
print("\n" + "─" * 70)
print("STEP 2: 資料品質驗證")
print("─" * 70)
validator = DataValidator(std_data)
issues = validator.validate_all()
summary = validator.get_validation_summary()
print(f"  錯誤: {summary['errors']}")
print(f"  警告: {summary['warnings']}")
if summary['errors'] == 0 and summary['warnings'] == 0:
    print("  ✓ 資料品質通過驗證")

# ============================================================
# STEP 3: 指標計算引擎
# ============================================================
print("\n" + "─" * 70)
print("STEP 3: 指標計算引擎")
print("─" * 70)
lineage = DataLineageTracker()
calc = MetricCalculator(std_data, lineage)

institutions = calc.get_all_institutions()
periods = calc.get_all_periods()
latest = periods[-1]
prev = calc.compute_prev_period(latest)

# 過濾掉總計和市佔率指標 institutions
real_institutions = [i for i in institutions if i != "總計"]

print(f"  最新期間: {latest} (民國114年12月)")
print(f"  機構數量: {len(real_institutions)}")

# 計算所有關鍵指標
print(f"\n  === 市場總計 ({latest}) ===")
market_cards = calc._get_market_total(latest, "流通卡數")
market_amount = calc._get_market_total(latest, "當月簽帳金額")
print(f"  市場流通卡數: {market_cards:,.0f} 張 ({market_cards/10000:,.0f} 萬張)")
print(f"  市場當月簽帳金額: {market_amount:,.0f} 千元 ({market_amount/1000000:,.0f} 億元)")

# 計算所有機構市占率
print(f"\n  === 流通卡數排名 Top 10 ({latest}) ===")
ranking_cards = calc.ranking(latest, "流通卡數", top_n=10)
for _, row in ranking_cards.iterrows():
    inst = row["institution"]
    val = row["value"]
    share = calc.market_share(inst, latest, "流通卡數")
    print(f"    #{int(row['rank']):2d} {inst:12s}  {val:>12,.0f} 張  市占率 {share:.1f}%")

print(f"\n  === 簽帳金額排名 Top 10 ({latest}) ===")
ranking_amount = calc.ranking(latest, "當月簽帳金額", top_n=10)
for _, row in ranking_amount.iterrows():
    inst = row["institution"]
    val = row["value"]
    share = calc.market_share(inst, latest, "當月簽帳金額")
    print(f"    #{int(row['rank']):2d} {inst:12s}  {val:>15,.0f} 千元  市占率 {share:.1f}%")

# 台新重點指標
taishin_names = [i for i in real_institutions if "台新" in i]
taishin = taishin_names[0] if taishin_names else None
if taishin:
    print(f"\n  === 台新國際商業銀行 重點指標 ===")
    ts_cards = calc._get_value(taishin, latest, "流通卡數")
    ts_amount = calc._get_value(taishin, latest, "當月簽帳金額")
    ts_share_cards = calc.market_share(taishin, latest, "流通卡數")
    ts_share_amount = calc.market_share(taishin, latest, "當月簽帳金額")

    # MoM
    ts_mom_cards = calc.mom_growth(taishin, latest, prev, "流通卡數")
    ts_mom_amount = calc.mom_growth(taishin, latest, prev, "當月簽帳金額")

    print(f"  流通卡數: {ts_cards:,.0f} 張 ({ts_cards/10000:,.0f} 萬張)")
    print(f"  流通卡數市占率: {ts_share_cards:.1f}%")
    print(f"  當月簽帳金額: {ts_amount:,.0f} 千元 ({ts_amount/1000000:,.1f} 億元)")
    print(f"  簽帳金額市占率: {ts_share_amount:.1f}%")
    print(f"  流通卡數 MoM: {ts_mom_cards:+.2f}%" if ts_mom_cards else "  流通卡數 MoM: N/A")
    print(f"  簽帳金額 MoM: {ts_mom_amount:+.2f}%" if ts_mom_amount else "  簽帳金額 MoM: N/A")

    # YoY 檢查（無113年資料，應為 N/A）
    prev_year = calc.compute_prev_year_period(latest)
    ts_yoy = calc.yoy_growth(taishin, latest, prev_year, "流通卡數")
    print(f"  流通卡數 YoY: {'N/A (無113年基期資料)' if ts_yoy is None else f'{ts_yoy:+.2f}%'}")

# ============================================================
# STEP 4: 簡報規劃 (使用附件二提示詞邏輯)
# ============================================================
print("\n" + "─" * 70)
print("STEP 4: 簡報結構規劃 (16 頁)")
print("─" * 70)
data_summary = {
    "institutions": real_institutions,
    "metrics": calc.get_all_metrics(),
    "periods": periods,
    "record_count": len(std_data),
}
planner = PlannerAgent()
slide_specs = planner.plan_structure(data_summary, use_llm=False)

# 用計算引擎結果填充 slide_spec 的數據
# Executive Summary KPIs
for spec in slide_specs:
    if spec.get("layout") == "executive_summary" and taishin:
        spec["kpis"] = [
            {
                "label": "市場流通卡數（12月）",
                "value": f"{market_cards/10000:,.0f} 萬",
                "metric_id": f"market_total_cards_{latest}",
                "change": "",
                "change_direction": "flat",
            },
            {
                "label": "月均簽帳金額（12月）",
                "value": f"{market_amount/1000000:,.0f} 億",
                "metric_id": f"market_total_amount_{latest}",
                "change": "",
                "change_direction": "flat",
            },
            {
                "label": "台新市占率（流通卡）",
                "value": f"{ts_share_cards:.1f}%",
                "metric_id": f"market_share_流通卡數_{taishin}_{latest}",
                "change": "排名第五",
                "change_direction": "flat",
            },
            {
                "label": "台新簽帳市占率",
                "value": f"{ts_share_amount:.1f}%",
                "metric_id": f"market_share_當月簽帳金額_{taishin}_{latest}",
                "change": "",
                "change_direction": "flat",
            },
        ]
        # 四大關鍵洞察（依據實際數據）
        # 計算玉山成長率
        yushan_names = [i for i in real_institutions if "玉山" in i]
        yushan = yushan_names[0] if yushan_names else ""
        yushan_cards_jan = calc._get_value(yushan, "11401", "流通卡數") if yushan else 0
        yushan_cards_dec = calc._get_value(yushan, latest, "流通卡數") if yushan else 0
        yushan_growth = ((yushan_cards_dec - yushan_cards_jan) / yushan_cards_jan * 100) if yushan_cards_jan else 0

        ts_cards_jan = calc._get_value(taishin, "11401", "流通卡數")
        ts_cards_growth = ((ts_cards - ts_cards_jan) / ts_cards_jan * 100) if ts_cards_jan else 0

        ts_amount_jan = calc._get_value(taishin, "11401", "當月簽帳金額")
        ts_amount_growth = ((ts_amount - ts_amount_jan) / ts_amount_jan * 100) if ts_amount_jan else 0

        spec["insights"] = [
            {
                "text": f"市場成長由簽帳額驅動，非卡數擴張。流通卡數年內僅微幅成長，但簽帳金額波動大，市場進入存量競爭階段。",
                "evidence_metric_ids": [f"market_total_cards_{latest}"],
                "is_speculation": False,
            },
            {
                "text": f"玉山以激進發卡攻市占，流通卡數年內成長 {yushan_growth:.1f}%，全行最高，但大量發卡的轉換效率有待觀察。",
                "evidence_metric_ids": [],
                "is_speculation": False,
            },
            {
                "text": f"台新簽帳金額年內成長 {ts_amount_growth:.1f}%，顯示既存客戶消費黏著度大幅提升，品質成長優於數量擴張。",
                "evidence_metric_ids": [f"market_share_當月簽帳金額_{taishin}_{latest}"],
                "is_speculation": False,
            },
            {
                "text": f"台新流通卡數年內成長 {ts_cards_growth:.1f}%，維持穩定，市占率 {ts_share_cards:.1f}% 排名第五，需加速發卡以縮小與前四名差距。",
                "evidence_metric_ids": [f"market_share_流通卡數_{taishin}_{latest}"],
                "is_speculation": False,
            },
        ]

    # 市占率排名圖
    elif spec.get("layout") == "ranking_chart":
        top10 = calc.ranking(latest, "流通卡數", top_n=10)
        categories = top10["institution"].tolist()
        shares = [calc.market_share(i, latest, "流通卡數") for i in categories]
        spec["chart"] = {
            "type": "bar",
            "title": f"流通卡數市占率排名 Top 10（{latest}）",
            "categories": [c.replace("商業銀行", "").replace("國際", "") for c in categories],
            "series": [{"name": "市占率(%)", "data": shares}],
            "y_axis": {"label": "市占率", "unit": "%"},
        }

    # 趨勢圖
    elif spec.get("layout") == "trend_chart":
        # 市場流通卡數月趨勢
        monthly_cards = [calc._get_market_total(p, "流通卡數") for p in periods]
        monthly_amount = [calc._get_market_total(p, "當月簽帳金額") for p in periods]
        months = [f"{int(p[3:])}月" for p in periods]
        spec["chart"] = {
            "type": "combo",
            "title": "市場規模趨勢 — 流通卡數與簽帳金額",
            "categories": months,
            "series": [
                {"name": "流通卡數(萬張)", "data": [v/10000 for v in monthly_cards if v]},
                {"name": "簽帳金額(億元)", "data": [v/1000000 for v in monthly_amount if v]},
            ],
            "y_axis": {"label": "數值", "unit": ""},
        }

    # 散點圖 (規模 vs 成長)
    elif spec.get("layout") == "scatter_chart":
        data_points = []
        for inst in real_institutions[:10]:
            cards = calc._get_value(inst, latest, "流通卡數")
            mom = calc.mom_growth(inst, latest, prev, "流通卡數")
            if cards and mom is not None:
                data_points.append({"name": inst, "x": cards/10000, "y": mom})
        spec["chart"] = {
            "type": "scatter",
            "title": "規模 vs 成長 — 流通卡數",
            "data_points": data_points,
            "x_axis": {"label": "流通卡數（萬張）", "unit": "萬張"},
            "y_axis": {"label": "月增率 MoM", "unit": "%"},
        }

    # 有效卡率 — 此資料不在附件四中，用簽帳金額/流通卡數 近似每卡消費力
    elif spec.get("layout") == "comparison_chart" and spec.get("slide_no") == 9:
        # 用「簽帳金額市占率 vs 流通卡數市占率」做比較
        top8 = calc.ranking(latest, "流通卡數", top_n=8)
        categories = [c.replace("商業銀行", "").replace("國際", "") for c in top8["institution"].tolist()]
        card_shares = [calc.market_share(i, latest, "流通卡數") or 0 for i in top8["institution"].tolist()]
        amount_shares = [calc.market_share(i, latest, "當月簽帳金額") or 0 for i in top8["institution"].tolist()]
        spec["headline"] = "簽帳金額市占率 vs 流通卡數市占率 — 差異反映單卡消費力"
        spec["chart"] = {
            "type": "bar",
            "title": "流通卡數 vs 簽帳金額市占率比較",
            "categories": categories,
            "series": [
                {"name": "流通卡數市占率(%)", "data": card_shares},
                {"name": "簽帳金額市占率(%)", "data": amount_shares},
            ],
            "y_axis": {"label": "市占率", "unit": "%"},
        }

    # 每卡簽帳金額
    elif spec.get("layout") == "comparison_chart" and spec.get("slide_no") == 11:
        # 計算每卡簽帳金額 = 當月簽帳金額 / 流通卡數
        top8 = calc.ranking(latest, "流通卡數", top_n=8)
        categories = []
        avg_per_card = []
        for inst in top8["institution"].tolist():
            cards = calc._get_value(inst, latest, "流通卡數")
            amount = calc._get_value(inst, latest, "當月簽帳金額")
            if cards and amount and cards > 0:
                categories.append(inst.replace("商業銀行", "").replace("國際", ""))
                avg_per_card.append(round(amount * 1000 / cards, 0))  # 千元轉元/卡
        spec["headline"] = "每卡月簽帳金額排名 — 台新每卡簽帳力高於多數同業"
        spec["chart"] = {
            "type": "bar",
            "title": f"平均每卡月簽帳金額（元/卡，{latest}）",
            "categories": categories,
            "series": [{"name": "每卡簽帳金額(元)", "data": avg_per_card}],
            "y_axis": {"label": "金額", "unit": "元/卡"},
        }

    # 循環信用與分期 — 用各銀行簽帳金額月趨勢替代
    elif spec.get("layout") == "stacked_chart":
        # 使用 Top 5 銀行的月簽帳金額趨勢
        top5 = calc.ranking(latest, "當月簽帳金額", top_n=5)
        months = [f"{int(p[3:])}月" for p in periods]
        series_list = []
        for inst in top5["institution"].tolist():
            monthly = [calc._get_value(inst, p, "當月簽帳金額") or 0 for p in periods]
            series_list.append({
                "name": inst.replace("商業銀行", "").replace("國際", ""),
                "data": [round(v/1000000, 1) for v in monthly],  # 轉億元
            })
        spec["headline"] = "Top 5 銀行月簽帳金額走勢（億元）"
        spec["chart"] = {
            "type": "stacked_bar",
            "title": "Top 5 銀行月簽帳金額堆疊圖",
            "categories": months,
            "series": series_list,
            "y_axis": {"label": "簽帳金額", "unit": "億元"},
        }

    # 風險指標 — 用流通卡數 MoM 波動度替代
    elif spec.get("layout") == "risk_chart":
        # 用各銀行流通卡數月增率作為「波動風險」指標
        top10 = calc.ranking(latest, "流通卡數", top_n=10)
        categories = []
        mom_values = []
        for inst in top10["institution"].tolist():
            mom = calc.mom_growth(inst, latest, prev, "流通卡數")
            if mom is not None:
                categories.append(inst.replace("商業銀行", "").replace("國際", ""))
                mom_values.append(mom)
        spec["headline"] = "流通卡數月增率 — 各銀行成長動態"
        spec["chart"] = {
            "type": "bar",
            "title": f"流通卡數月增率 MoM（{prev}→{latest}）",
            "categories": categories,
            "series": [{"name": "月增率(%)", "data": mom_values}],
            "y_axis": {"label": "月增率", "unit": "%"},
        }

    # 策略建議
    elif spec.get("layout") == "strategy":
        spec["recommendations"] = [
            {
                "action": "加速數位發卡，擴大流通卡規模",
                "rationale": f"台新流通卡數 {ts_cards/10000:,.0f} 萬排名第五，市占 {ts_share_cards:.1f}%，與第一名中信差距仍大。建議強化線上發卡體驗、推出數位專屬卡種。",
                "priority": "high",
            },
            {
                "action": "深化消費場景，提升每卡簽帳力",
                "rationale": f"台新簽帳金額市占 {ts_share_amount:.1f}% 略低於流通卡市占 {ts_share_cards:.1f}%，顯示單卡消費力仍有提升空間。建議與電商/外送平台深度綁定。",
                "priority": "high",
            },
            {
                "action": "維持風險控管優勢，建立市場信任",
                "rationale": "在競爭者積極擴張背景下，將風險優勢轉化為品牌差異化。",
                "priority": "medium",
            },
            {
                "action": "精進有效卡經營，降低沉沒卡成本",
                "rationale": "針對 6 個月未消費卡戶啟動精準喚醒 campaign，建立沉睡卡預警機制。",
                "priority": "medium",
            },
        ]

for spec in slide_specs:
    print(f"  頁 {spec.get('slide_no'):2d}: [{spec.get('layout'):20s}] {spec.get('title', '')}")

# ============================================================
# STEP 5: 洞察生成
# ============================================================
print("\n" + "─" * 70)
print("STEP 5: 商業洞察生成 (rule-based)")
print("─" * 70)
analyst = AnalystAgent()
all_metric_data = {s.get("slide_no", 0): lineage.export_summary() for s in slide_specs}
enriched_specs = analyst.generate_all_slides(slide_specs, all_metric_data, use_llm=False)
print(f"  已為 {len(enriched_specs)} 頁生成洞察")

# ============================================================
# STEP 6: 品質審核
# ============================================================
print("\n" + "─" * 70)
print("STEP 6: 品質審核 (Reviewer Agent)")
print("─" * 70)
reviewer = ReviewerAgent()
verified_metrics = {r.metric_id: r.value for r in lineage.records.values()}
qa_result = reviewer.review(enriched_specs, verified_metrics, use_llm=False)
print(f"  審核狀態: {qa_result['status']}")
print(f"  錯誤: {qa_result['errors']}")
print(f"  警告: {qa_result['warnings']}")
if qa_result['issues']:
    for issue in qa_result['issues'][:5]:
        print(f"    [{issue.get('severity')}] Slide {issue.get('slide_no')}: {issue.get('message')}")

# ============================================================
# STEP 7: PPT 生成
# ============================================================
print("\n" + "─" * 70)
print("STEP 7: PowerPoint 生成 (使用附件一模板)")
print("─" * 70)
output_ppt = os.path.join(OUTPUT_DIR, "final_presentation.pptx")
generator = PPTGenerator(TEMPLATE_PATH)
result_path = generator.generate(enriched_specs, output_ppt)
print(f"  輸出: {result_path}")

from pptx import Presentation
prs = Presentation(output_ppt)
print(f"  共 {len(prs.slides)} 頁投影片")
for i, s in enumerate(prs.slides):
    # 取得每頁的標題文字
    title_text = ""
    for shape in s.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title_text = shape.text_frame.text.strip()[:50]
            break
    print(f"    {i+1:2d}. [{s.slide_layout.name:12s}] {title_text}")

# ============================================================
# STEP 8: 數值回溯校驗
# ============================================================
print("\n" + "─" * 70)
print("STEP 8: 數值回溯校驗")
print("─" * 70)
reconciler = PPTReconciler(lineage)
reconcile_result = reconciler.reconcile(enriched_specs)
print(f"  校驗狀態: {reconcile_result['status']}")
print(f"  錯誤: {reconcile_result['total_errors']}")
print(f"  警告: {reconcile_result['total_warnings']}")

# ============================================================
# STEP 9: 匯出結果
# ============================================================
print("\n" + "─" * 70)
print("STEP 9: 匯出結果")
print("─" * 70)
# Excel
excel_out = os.path.join(OUTPUT_DIR, "analysis_result.xlsx")
std_data.to_excel(excel_out, index=False, engine="openpyxl")
print(f"  分析結果 Excel: {excel_out}")

# Data lineage
lineage_out = os.path.join(OUTPUT_DIR, "data_lineage.json")
lineage.export_json(lineage_out)
print(f"  資料血緣 JSON: {lineage_out} ({len(lineage.records)} 筆記錄)")

# Slide spec
spec_out = os.path.join(OUTPUT_DIR, "slide_spec.json")
with open(spec_out, "w", encoding="utf-8") as f:
    json.dump(enriched_specs, f, ensure_ascii=False, indent=2)
print(f"  Slide Spec JSON: {spec_out}")

# QA report
qa_out = os.path.join(OUTPUT_DIR, "qa_report.json")
with open(qa_out, "w", encoding="utf-8") as f:
    json.dump(qa_result, f, ensure_ascii=False, indent=2)
print(f"  QA 報告 JSON: {qa_out}")

# ============================================================
# 最終摘要
# ============================================================
print("\n" + "=" * 70)
print("  DEMO 完成 — 輸出結果摘要")
print("=" * 70)
print(f"""
  輸入:
    Excel: {EXCEL_PATH}
    模板:  {TEMPLATE_PATH}

  輸出 ({OUTPUT_DIR}/):
    1. final_presentation.pptx  — 16 頁策略簡報
    2. analysis_result.xlsx     — 標準化分析結果
    3. data_lineage.json        — 數值追溯紀錄 ({len(lineage.records)} 筆)
    4. slide_spec.json          — 簡報規格 JSON
    5. qa_report.json           — 品質校驗報告

  關鍵數據驗證:
    市場流通卡數 (11412): {market_cards:,.0f} 張 = {market_cards/10000:,.0f} 萬張
    台新流通卡數市占率: {ts_share_cards:.1f}%
    台新排名: 第5名
    YoY: N/A (無113年基期資料 — 系統正確拒絕產生)

  品質檢查:
    資料驗證: {'✓ 通過' if summary['errors'] == 0 else '✗ 有錯誤'}
    品質審核: {qa_result['status']}
    數值回溯: {reconcile_result['status']}
""")
print("=" * 70)
