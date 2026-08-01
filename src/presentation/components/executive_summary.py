"""Executive Summary 頁面元件 — KPI 卡片 + 洞察。
支援智慧版面自適應：KPI 數量少時放大居中，洞察多時自動縮小字體。
支援條件式格式：負成長紅色數字。
"""

from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.slide import Slide


def fill_executive_summary(slide: Slide, spec: dict):
    """填充 Executive Summary 頁 — KPI 卡片 + 洞察（自適應版面）。"""
    kpis = spec.get("kpis", [])
    kpi_count = min(len(kpis), 4)

    # --- 智慧版面自適應：根據 KPI 數量調整卡片大小和位置 ---
    if kpi_count <= 2:
        # 少量 KPI：放大居中
        card_width = Cm(12.0)
        total_width = card_width * kpi_count + Cm(2.0) * (kpi_count - 1) if kpi_count > 1 else card_width
        start_x = Cm((33.9 - total_width / 914400 * 2.54) / 2) if kpi_count == 1 else Cm(5.0)
        gap = Cm(14.0) if kpi_count > 1 else Cm(0)
        label_size = Pt(14)
        value_size = Pt(28)
        change_size = Pt(12)
    elif kpi_count == 3:
        card_width = Cm(9.5)
        start_x = Cm(1.5)
        gap = Cm(10.3)
        label_size = Pt(12)
        value_size = Pt(24)
        change_size = Pt(11)
    else:
        # 4 個 KPI：標準布局
        card_width = Cm(7.5)
        start_x = Cm(1.0)
        gap = Cm(7.8)
        label_size = Pt(12)
        value_size = Pt(24)
        change_size = Pt(11)

    y = Cm(4.0)

    for i, kpi in enumerate(kpis[:4]):
        x = start_x + Cm(i * (gap / 914400 * 2.54)) if kpi_count > 2 else (start_x + Cm(i * 14.0) if kpi_count == 2 else start_x)
        # 簡化計算：使用固定間距
        if kpi_count <= 2:
            x = Cm(5.0 + i * 14.0) if kpi_count == 2 else Cm(11.0)
        elif kpi_count == 3:
            x = Cm(1.5 + i * 10.3)
        else:
            x = Cm(1.0 + i * 7.8)

        # 標籤
        txBox = slide.shapes.add_textbox(x, y, card_width, Cm(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = kpi.get("label", "")
        p.font.size = label_size
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        # 數值
        txBox2 = slide.shapes.add_textbox(x, y + Cm(0.9), card_width, Cm(1.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = kpi.get("value", "—")
        p2.font.size = value_size
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

        # 變化 — 條件式格式：負成長紅色
        change = kpi.get("change", "")
        if change:
            txBox3 = slide.shapes.add_textbox(x, y + Cm(2.5), card_width, Cm(0.7))
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = change
            p3.font.size = change_size
            direction = kpi.get("change_direction", "flat")
            if direction == "up":
                p3.font.color.rgb = RGBColor(0x33, 0x99, 0x33)  # 綠色
            elif direction == "down":
                p3.font.color.rgb = RGBColor(0xCC, 0x33, 0x33)  # 紅色
                p3.font.bold = True  # 負成長加粗醒目
            else:
                p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- 洞察區域（自適應字體大小） ---
    insights = spec.get("insights", [])
    if insights:
        insight_y = Cm(8.0)
        # 洞察標題
        txH = slide.shapes.add_textbox(Cm(1.0), insight_y, Cm(31.2), Cm(0.8))
        tfH = txH.text_frame
        pH = tfH.paragraphs[0]
        pH.text = "關鍵洞察"
        pH.font.size = Pt(14)
        pH.font.bold = True
        pH.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

        insight_y += Cm(1.0)

        # 智慧自適應：洞察超過 4 條時縮小字體和行距
        insight_count = min(len(insights), 6)
        if insight_count <= 3:
            font_size = Pt(12)
            line_height = Cm(1.5)
        elif insight_count <= 4:
            font_size = Pt(11)
            line_height = Cm(1.4)
        else:
            font_size = Pt(10)
            line_height = Cm(1.2)

        for insight in insights[:insight_count]:
            text = insight.get("text", "") if isinstance(insight, dict) else str(insight)
            txBox = slide.shapes.add_textbox(Cm(1.5), insight_y, Cm(30.0), line_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {text}"
            p.font.size = font_size
            insight_y += line_height
