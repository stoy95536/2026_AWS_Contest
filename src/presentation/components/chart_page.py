"""圖表頁面元件 — 使用 python-pptx 原生圖表填充內容區域。"""

from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide

from ..chart_factory import ChartFactory


def fill_chart(slide: Slide, spec: dict, chart_factory: ChartFactory):
    """
    填充含圖表的頁面 — 使用 python-pptx 原生圖表。

    Args:
        slide: 目標投影片
        spec: slide_spec dict
        chart_factory: ChartFactory 實例
    """
    chart_spec = spec.get("chart", {})
    chart_type = chart_spec.get("type", "bar")
    categories = chart_spec.get("categories", [])
    series_list = chart_spec.get("series", [])
    data_points = chart_spec.get("data_points", [])

    # 圖表位置 (在 1_標題及內容 的 body 區域內)
    position = {
        "left": Cm(1.5),
        "top": Cm(4.5),
        "width": Cm(30.0),
        "height": Cm(10.0),
    }

    chart_created = False

    # 散佈圖使用 data_points
    if chart_type == "scatter" and data_points:
        chart_factory.create_scatter_chart(
            slide, data_points,
            title=chart_spec.get("title", ""),
            x_label=chart_spec.get("x_axis", {}).get("label", ""),
            y_label=chart_spec.get("y_axis", {}).get("label", ""),
            position=position,
        )
        chart_created = True

    # 其他圖表使用 categories + series
    elif series_list and categories:
        series_data = {}
        for s in series_list:
            if "data" in s and s["data"]:
                series_data[s.get("name", "Series")] = s["data"]

        if series_data:
            y_axis_label = chart_spec.get("y_axis", {}).get("label", "")
            y_axis_unit = chart_spec.get("y_axis", {}).get("unit", "")

            if chart_type == "bar":
                chart_factory.create_bar_chart(
                    slide, categories, series_data,
                    title=chart_spec.get("title", ""),
                    position=position,
                    y_axis_label=y_axis_label,
                    y_axis_unit=y_axis_unit,
                )
                chart_created = True
            elif chart_type == "line":
                chart_factory.create_line_chart(
                    slide, categories, series_data,
                    title=chart_spec.get("title", ""),
                    position=position,
                    y_axis_label=y_axis_label,
                    y_axis_unit=y_axis_unit,
                )
                chart_created = True
            elif chart_type == "combo":
                keys = list(series_data.keys())
                bar_data = {keys[0]: series_data[keys[0]]} if keys else {}
                line_data = {k: series_data[k] for k in keys[1:]} if len(keys) > 1 else {}
                chart_factory.create_combo_chart(
                    slide, categories, bar_data, line_data,
                    title=chart_spec.get("title", ""),
                    position=position,
                    y_axis_label=y_axis_label,
                    y_axis_unit=y_axis_unit,
                )
                chart_created = True
            elif chart_type == "stacked_bar":
                chart_factory.create_stacked_bar_chart(
                    slide, categories, series_data,
                    title=chart_spec.get("title", ""),
                    position=position,
                    y_axis_label=y_axis_label,
                    y_axis_unit=y_axis_unit,
                )
                chart_created = True
            elif chart_type == "pie":
                first_series = series_list[0] if series_list else {}
                values = first_series.get("data", [])
                if values:
                    chart_factory.create_pie_chart(
                        slide, categories, values,
                        title=chart_spec.get("title", ""),
                        position=position,
                    )
                    chart_created = True

    # 如果沒有生成圖表，顯示提示
    if not chart_created:
        txBox = slide.shapes.add_textbox(Cm(5), Cm(7), Cm(23), Cm(3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[圖表區域 — {chart_spec.get('title', '待補充資料')}]"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER

    # 洞察文字 (圖表下方)
    insights = spec.get("insights", [])
    if insights:
        insight_y = Cm(15.0)
        for ins in insights[:2]:
            text = ins.get("text", "") if isinstance(ins, dict) else str(ins)
            is_spec = ins.get("is_speculation", False) if isinstance(ins, dict) else False
            txBox = slide.shapes.add_textbox(Cm(1.5), insight_y, Cm(30.0), Cm(1.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if is_spec:
                # 推測性內容：橘色斜體，讓評審一眼區分 LLM 推論
                p.text = f"◆ （推測）{text}"
                p.font.size = Pt(11)
                p.font.italic = True
                p.font.color.rgb = RGBColor(0xFF, 0x66, 0x00)  # 台新橘
            else:
                p.text = f"◆ {text}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            insight_y += Cm(1.3)
