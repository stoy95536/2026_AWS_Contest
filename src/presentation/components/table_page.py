"""表格頁面元件 — 含 header 背景色、zebra striping、邊框樣式。"""

from lxml import etree

from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.slide import Slide


def fill_table(slide: Slide, spec: dict):
    """
    填充含原生表格的頁面。
    包含 header 背景色（台新深藍）、交替行底色 (zebra striping)、邊框樣式。
    """
    table_spec = spec.get("table", {})
    headers = table_spec.get("headers", [])
    rows = table_spec.get("rows", [])

    if not headers or not rows:
        return

    n_rows = min(len(rows) + 1, 15)
    n_cols = len(headers)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Cm(1.5), Cm(4.5), Cm(30.0), Cm(11.0),
    )
    table = tbl_shape.table

    col_width = Cm(30.0 / n_cols)
    for col_idx in range(n_cols):
        table.columns[col_idx].width = int(col_width)

    # Header 行
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        _set_cell_fill(cell, RGBColor(0x00, 0x33, 0x66))
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(9)
            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_border(cell, RGBColor(0xFF, 0xFF, 0xFF), Pt(0.5))

    # Data rows
    even_row_color = RGBColor(0xE8, 0xF0, 0xF8)
    odd_row_color = RGBColor(0xFF, 0xFF, 0xFF)
    border_color = RGBColor(0xCC, 0xCC, 0xCC)

    for row_idx, row_data in enumerate(rows[:n_rows - 1]):
        bg_color = even_row_color if row_idx % 2 == 0 else odd_row_color
        for col_idx, value in enumerate(row_data):
            if col_idx < n_cols:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value) if value is not None else ""
                _set_cell_fill(cell, bg_color)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(8)
                    para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                _set_cell_border(cell, border_color, Pt(0.5))


def _set_cell_fill(cell, color: RGBColor):
    """設定表格儲存格的背景填充色。"""
    tc_pr = cell._tc.get_or_add_tcPr()
    for old_fill in tc_pr.findall(qn("a:solidFill")):
        tc_pr.remove(old_fill)
    solid_fill = etree.SubElement(tc_pr, qn("a:solidFill"))
    srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb_clr.set("val", str(color))


def _set_cell_border(cell, color: RGBColor, width):
    """設定表格儲存格的四邊邊框。"""
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    for border_name in borders:
        for old_ln in tc_pr.findall(qn(border_name)):
            tc_pr.remove(old_ln)
        ln = etree.SubElement(tc_pr, qn(border_name))
        ln.set("w", str(int(width)))
        solid_fill = etree.SubElement(ln, qn("a:solidFill"))
        srgb_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
        srgb_clr.set("val", str(color))
