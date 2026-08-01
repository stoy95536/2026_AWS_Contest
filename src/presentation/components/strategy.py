"""策略建議頁面元件。"""

from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.slide import Slide


def fill_strategy(slide: Slide, spec: dict):
    """填充策略建議頁。"""
    recommendations = spec.get("recommendations", [])
    y = Cm(4.0)

    for i, rec in enumerate(recommendations[:4]):
        action = rec.get("action", "")
        rationale = rec.get("rationale", "")
        priority = rec.get("priority", "medium")

        # 編號
        txNum = slide.shapes.add_textbox(Cm(1.5), y, Cm(2.0), Cm(1.5))
        tfNum = txNum.text_frame
        pNum = tfNum.paragraphs[0]
        pNum.text = f"0{i+1}"
        pNum.font.size = Pt(24)
        pNum.font.bold = True
        pNum.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

        # 標題
        txTitle = slide.shapes.add_textbox(Cm(4.0), y, Cm(27.0), Cm(0.8))
        tfTitle = txTitle.text_frame
        pTitle = tfTitle.paragraphs[0]
        pTitle.text = action
        pTitle.font.size = Pt(13)
        pTitle.font.bold = True
        pTitle.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

        # 理由
        if rationale:
            txR = slide.shapes.add_textbox(Cm(4.0), y + Cm(0.9), Cm(27.0), Cm(1.5))
            tfR = txR.text_frame
            tfR.word_wrap = True
            pR = tfR.paragraphs[0]
            pR.text = rationale
            pR.font.size = Pt(10)
            pR.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        y += Cm(3.0)
