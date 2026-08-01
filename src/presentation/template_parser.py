"""
PowerPoint 模板解析器
支援自動偵測任意模板的 layout 配置，不再寫死名稱。
透過分析每個 layout 的 placeholder 特徵，自動對應到：
  cover / content / chapter / thank_you 四種頁面類型。
"""

import os
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# 向後相容：保留舊常數供外部 import 使用
LAYOUT_COVER = "2_標題投影片"
LAYOUT_CONTENT = "1_標題及內容"
LAYOUT_CHAPTER = "2_章節標題"
LAYOUT_THANKYOU = "3_標題投影片"
LAYOUT_TITLE_BODY = "標題投影片"
LAYOUT_CENTER = "1_標題投影片"


@dataclass
class TemplateStyle:
    """模板樣式設定。"""
    slide_width: int = 12192000  # 33.9 cm (16:9)
    slide_height: int = 6858000  # 19.1 cm
    title_font: str = "微軟正黑體"
    title_size: int = 28
    body_font: str = "微軟正黑體"
    body_size: int = 14
    primary_color: str = "003366"  # 台新深藍
    secondary_color: str = "0066CC"
    accent_color: str = "FF6600"
    # 內容頁的座標（自動偵測後會更新）
    content_title_left: float = 1.0   # cm
    content_title_top: float = 0.7
    content_title_width: float = 31.3
    content_title_height: float = 2.1
    content_body_left: float = 1.0
    content_body_top: float = 3.5
    content_body_width: float = 31.2
    content_body_height: float = 13.6
    slide_number_left: float = 31.7
    slide_number_top: float = 17.8


@dataclass
class LayoutInfo:
    """單一 layout 的偵測結果。"""
    name: str
    index: int
    has_title: bool = False
    has_body: bool = False
    has_slide_number: bool = False
    has_subtitle: bool = False
    placeholder_count: int = 0
    title_top_ratio: float = 0.0  # 標題 top 位置佔投影片高度比例


class TemplateParser:
    """
    解析 PowerPoint 模板，自動偵測 layout 對應。
    支援任意模板，不需要手動指定 layout 名稱。
    """

    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        self.style = TemplateStyle()
        self._layout_map: dict[str, int] = {}
        self._layout_infos: list[LayoutInfo] = []
        # 自動偵測結果：page_type → layout_name
        self._type_to_layout: dict[str, str] = {}

        if template_path and os.path.exists(template_path):
            self._parse_template()
            self._auto_detect_layouts()
        else:
            # 無模板時使用預設常數
            self._type_to_layout = {
                "cover": LAYOUT_COVER,
                "content": LAYOUT_CONTENT,
                "chapter": LAYOUT_CHAPTER,
                "thank_you": LAYOUT_THANKYOU,
            }

    def _parse_template(self):
        """解析模板檔案，建立 layout 名稱索引和特徵資訊。"""
        prs = Presentation(self.template_path)
        self.style.slide_width = prs.slide_width
        self.style.slide_height = prs.slide_height
        slide_height = prs.slide_height

        for i, layout in enumerate(prs.slide_layouts):
            self._layout_map[layout.name] = i

            # 分析 placeholder 特徵
            info = LayoutInfo(name=layout.name, index=i)
            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx
                ph_type = ph.placeholder_format.type
                info.placeholder_count += 1

                # idx=0 通常是 TITLE
                if idx == 0:
                    info.has_title = True
                    # 計算標題 top 位置佔比（判斷是否置中/偏下）
                    if ph.top and slide_height:
                        info.title_top_ratio = ph.top / slide_height
                # idx=1 通常是 BODY
                elif idx == 1:
                    info.has_body = True
                # idx=12 或 idx=10 通常是 SLIDE_NUMBER
                elif idx in (10, 12):
                    info.has_slide_number = True
                # idx=13 或特定類型是 SUBTITLE
                elif idx == 1 and str(ph_type) == "SUBTITLE (4)":
                    info.has_subtitle = True

            self._layout_infos.append(info)

    def _auto_detect_layouts(self):
        """
        自動偵測 layout 用途。
        規則：
          content  = 有 title + 有 body（或 placeholder >= 3）→ 內容頁
          chapter  = 有 title + 無 body + 標題偏下（top > 40%）→ 章節頁
          cover    = 有 title + 無 body + 標題偏上/中 → 封面
          thank_you = 最後一個 title-only layout（跟 cover 類似但排後面）
        """
        content_candidates = []
        title_only_candidates = []
        chapter_candidates = []

        for info in self._layout_infos:
            if info.has_title and info.has_body:
                content_candidates.append(info)
            elif info.has_title and not info.has_body:
                if info.title_top_ratio > 0.4:
                    chapter_candidates.append(info)
                else:
                    title_only_candidates.append(info)

        # content：優先選有 slide_number 的，或 placeholder 最多的
        if content_candidates:
            # 優先有頁碼的
            with_number = [c for c in content_candidates if c.has_slide_number]
            chosen = with_number[0] if with_number else content_candidates[0]
            self._type_to_layout["content"] = chosen.name
        else:
            # 沒有就用第一個 layout
            self._type_to_layout["content"] = self._layout_infos[0].name if self._layout_infos else ""

        # chapter：標題偏下的
        if chapter_candidates:
            self._type_to_layout["chapter"] = chapter_candidates[0].name
        else:
            # fallback 到 title_only 的第二個（如果有）
            fallback = title_only_candidates[1] if len(title_only_candidates) > 1 else None
            self._type_to_layout["chapter"] = fallback.name if fallback else self._type_to_layout["content"]

        # cover 和 thank_you：從 title_only 候選中分配
        if title_only_candidates:
            self._type_to_layout["cover"] = title_only_candidates[0].name
            # thank_you 用最後一個 title_only（或跟 cover 不同的）
            if len(title_only_candidates) >= 2:
                self._type_to_layout["thank_you"] = title_only_candidates[-1].name
            else:
                self._type_to_layout["thank_you"] = title_only_candidates[0].name
        else:
            # 全部 fallback
            self._type_to_layout["cover"] = self._type_to_layout["content"]
            self._type_to_layout["thank_you"] = self._type_to_layout["content"]

        # 也嘗試用名稱關鍵字做精確匹配（如果名字包含特定文字，優先使用）
        self._refine_by_name_keywords()

    def _refine_by_name_keywords(self):
        """用 layout 名稱關鍵字精確匹配，覆蓋自動偵測結果。"""
        keyword_map = {
            "content": ["內容", "content", "標題及內容"],
            "cover": ["封面", "cover", "標題投影片"],
            "chapter": ["章節", "chapter", "section"],
            "thank_you": ["結束", "thank", "ending", "end"],
        }

        for info in self._layout_infos:
            name_lower = info.name.lower()
            for page_type, keywords in keyword_map.items():
                for kw in keywords:
                    if kw.lower() in name_lower:
                        # 特殊處理：「標題及內容」精確匹配 content
                        if page_type == "content" and "內容" in info.name:
                            self._type_to_layout["content"] = info.name
                        elif page_type == "chapter" and ("章節" in info.name or "section" in name_lower):
                            self._type_to_layout["chapter"] = info.name
                        break

    def get_style(self) -> TemplateStyle:
        """取得模板樣式。"""
        return self.style

    def get_layout_index(self, layout_name: str) -> int:
        """取得指定 layout 的索引。"""
        return self._layout_map.get(layout_name, 0)

    def get_detected_mapping(self) -> dict[str, str]:
        """取得自動偵測的 page_type → layout_name 對應（方便 debug）。"""
        return dict(self._type_to_layout)

    def create_presentation(self) -> Presentation:
        """
        建立簡報（基於模板）。
        會清除模板中既有的範例投影片，只保留 layout 定義。
        """
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
            # 清除模板中既有的範例投影片
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            return prs
        else:
            prs = Presentation()
            prs.slide_width = self.style.slide_width
            prs.slide_height = self.style.slide_height
            return prs

    def get_layout_for_page_type(self, page_type: str) -> str:
        """
        根據頁面類型回傳對應的 layout 名稱。
        使用自動偵測結果，而非寫死的常數。

        Page types: cover, content, chapter, thank_you
        """
        return self._type_to_layout.get(page_type, self._type_to_layout.get("content", ""))
