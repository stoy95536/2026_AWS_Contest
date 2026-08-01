"""
圖表工廠
建立 PowerPoint 原生可編輯圖表。
所有圖表為原生物件，不以圖片嵌入。
"""

import copy

from lxml import etree

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from pptx.slide import Slide


# 台新品牌配色
TAISHIN_COLORS = [
    RGBColor(0x00, 0x33, 0x66),  # 深藍
    RGBColor(0x00, 0x66, 0xCC),  # 亮藍
    RGBColor(0xFF, 0x66, 0x00),  # 橘色
    RGBColor(0x33, 0x99, 0x33),  # 綠色
    RGBColor(0xCC, 0x33, 0x33),  # 紅色
    RGBColor(0x99, 0x66, 0xCC),  # 紫色
    RGBColor(0xFF, 0xCC, 0x00),  # 黃色
    RGBColor(0x66, 0x99, 0x99),  # 青色
]


class ChartFactory:
    """建立 PowerPoint 原生可編輯圖表。"""

    def __init__(self):
        self.colors = TAISHIN_COLORS

    def create_bar_chart(
        self,
        slide: Slide,
        categories: list[str],
        series_data: dict[str, list[float]],
        title: str = "",
        position: dict = None,
        y_axis_label: str = "",
        y_axis_unit: str = "",
        highlight_first: bool = True,
        highlight_institution: str = "",
    ):
        """
        建立橫條/直條圖。
        支援條件式格式：排名第一金色、目標機構深藍高亮。

        Args:
            slide: 目標投影片
            categories: X 軸類別
            series_data: {系列名稱: [值...]}
            title: 圖表標題
            position: 位置 dict (left, top, width, height)
            y_axis_label: Y 軸標籤
            y_axis_unit: Y 軸單位
            highlight_first: 是否將排名第一的長條標為金色
            highlight_institution: 目標機構名稱（深藍高亮）
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = CategoryChartData()
        chart_data.categories = categories

        for name, values in series_data.items():
            chart_data.add_series(name, values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # 設定顏色 — 條件式格式
        for i, series in enumerate(chart.series):
            base_color = self.colors[i % len(self.colors)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = base_color

            # 條件式格式：逐點上色（排名第一金色、目標機構深藍）
            if len(series_data) == 1:  # 單系列才做逐點高亮
                for pt_idx, cat in enumerate(categories):
                    point = series.points[pt_idx]
                    if highlight_first and pt_idx == 0:
                        # 排名第一：金色
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0x00)
                    elif highlight_institution and highlight_institution in cat:
                        # 目標機構：台新深藍
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

        # Y 軸設定
        self._set_value_axis_title(chart.value_axis, y_axis_label, y_axis_unit)

        return chart_frame

    def create_line_chart(
        self,
        slide: Slide,
        categories: list[str],
        series_data: dict[str, list[float]],
        title: str = "",
        position: dict = None,
        y_axis_label: str = "",
        y_axis_unit: str = "",
    ):
        """
        建立折線圖。

        Args:
            slide: 目標投影片
            categories: X 軸類別
            series_data: {系列名稱: [值...]}
            title: 圖表標題
            position: 位置 dict (left, top, width, height)
            y_axis_label: Y 軸標籤
            y_axis_unit: Y 軸單位
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = CategoryChartData()
        chart_data.categories = categories

        for name, values in series_data.items():
            chart_data.add_series(name, values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        for i, series in enumerate(chart.series):
            series.format.line.color.rgb = self.colors[i % len(self.colors)]
            series.format.line.width = Pt(2.5)

        # Y 軸設定
        self._set_value_axis_title(chart.value_axis, y_axis_label, y_axis_unit)

        return chart_frame

    def create_combo_chart(
        self,
        slide: Slide,
        categories: list[str],
        bar_series: dict[str, list[float]],
        line_series: dict[str, list[float]],
        title: str = "",
        position: dict = None,
        y_axis_label: str = "",
        y_axis_unit: str = "",
    ):
        """
        建立組合圖（直條 + 折線）。
        透過操作底層 XML，將折線系列從 barChart plot 移至獨立的 lineChart plot，
        實現真正的組合圖效果。

        Args:
            slide: 目標投影片
            categories: X 軸類別
            bar_series: 直條系列 {系列名稱: [值...]}
            line_series: 折線系列 {系列名稱: [值...]}
            title: 圖表標題
            position: 位置 dict (left, top, width, height)
            y_axis_label: Y 軸標籤
            y_axis_unit: Y 軸單位
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = CategoryChartData()
        chart_data.categories = categories

        # 先加直條系列
        for name, values in bar_series.items():
            chart_data.add_series(name, values)
        # 再加折線系列
        for name, values in line_series.items():
            chart_data.add_series(name, values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # --- 透過 XML 操作將折線系列移至獨立 lineChart plot ---
        bar_count = len(bar_series)
        if line_series and bar_count < len(chart.series):
            plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
            bar_chart_elem = plot_area.find(qn("c:barChart"))

            if bar_chart_elem is not None:
                # 建立 lineChart 元素
                line_chart_elem = etree.SubElement(plot_area, qn("c:lineChart"))
                # 設定 grouping 為 standard
                grouping = etree.SubElement(line_chart_elem, qn("c:grouping"))
                grouping.set("val", "standard")

                # 從 barChart 中取出折線系列，移至 lineChart
                all_ser = bar_chart_elem.findall(qn("c:ser"))
                for i in range(bar_count, bar_count + len(line_series)):
                    if i < len(all_ser):
                        ser_elem = all_ser[i]
                        bar_chart_elem.remove(ser_elem)
                        line_chart_elem.append(ser_elem)

                # 為 lineChart 加入 marker 設定（顯示標記點）
                marker = etree.SubElement(line_chart_elem, qn("c:marker"))
                marker_val = etree.SubElement(marker, qn("c:val"))
                marker_val.text = "1"

                # 設定折線系列使用副Y軸 (axId 對應)
                # 複製主軸 axId 引用到 lineChart
                for ax_id in bar_chart_elem.findall(qn("c:axId")):
                    new_ax_id = copy.deepcopy(ax_id)
                    line_chart_elem.append(new_ax_id)

        # 設定直條系列顏色
        for i, series in enumerate(chart.series):
            if i < bar_count:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self.colors[i % len(self.colors)]
            else:
                # 折線系列設定線條顏色
                color_idx = i % len(self.colors)
                series.format.line.color.rgb = self.colors[color_idx]
                series.format.line.width = Pt(2.5)

        # Y 軸設定
        self._set_value_axis_title(chart.value_axis, y_axis_label, y_axis_unit)

        return chart_frame

    def create_scatter_chart(
        self,
        slide: Slide,
        data_points: list[dict],
        title: str = "",
        x_label: str = "",
        y_label: str = "",
        position: dict = None,
    ):
        """
        建立散佈圖，並為每個資料點加上名稱標籤。

        Args:
            data_points: [{"name": str, "x": float, "y": float}, ...]
            title: 圖表標題
            x_label: X 軸標籤
            y_label: Y 軸標籤
            position: 位置 dict (left, top, width, height)
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = XyChartData()

        series = chart_data.add_series("銀行")
        for point in data_points:
            series.add_data_point(point["x"], point["y"])

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # X 軸標籤
        if x_label:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = x_label

        # Y 軸標籤
        if y_label:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_label

        # 將 X 軸刻度移到底部（Y 軸在最小值處交叉），避免與資料點重疊
        # python-pptx 的 crosses 屬性在散佈圖上不一定生效，直接操作 XML
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        # 找到 Y 軸 (valAx)，設定它在 X 軸最小值處交叉
        for val_ax in plot_area.findall(qn("c:valAx")):
            # 找到 crossesAt 或 crosses 元素
            crosses = val_ax.find(qn("c:crosses"))
            if crosses is not None:
                crosses.set("val", "min")
            else:
                crosses = etree.SubElement(val_ax, qn("c:crosses"))
                crosses.set("val", "min")

        # --- 透過 XML 為每個資料點加上名稱標籤 ---
        self._add_scatter_data_labels(chart, data_points)

        return chart_frame

    def _add_scatter_data_labels(self, chart, data_points: list[dict]):
        """
        為散佈圖的每個資料點加入自訂文字標籤（銀行名稱）。
        標籤放置於資料點上方，避免與軸刻度重疊。
        """
        plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        scatter_chart = plot_area.find(qn("c:scatterChart"))
        if scatter_chart is None:
            return

        ser_elem = scatter_chart.find(qn("c:ser"))
        if ser_elem is None:
            return

        # 建立 dLbls 容器
        d_lbls = etree.SubElement(ser_elem, qn("c:dLbls"))

        for idx, point in enumerate(data_points):
            name = point.get("name", "")
            if not name:
                continue

            d_lbl = etree.SubElement(d_lbls, qn("c:dLbl"))

            # 指定 data point index
            idx_elem = etree.SubElement(d_lbl, qn("c:idx"))
            idx_elem.set("val", str(idx))

            # 標籤位置設為上方 (t = top)，避免跟軸線重疊
            d_lbl_pos = etree.SubElement(d_lbl, qn("c:dLblPos"))
            d_lbl_pos.set("val", "t")

            # 使用 tx (rich text) 設定自訂標籤文字
            tx = etree.SubElement(d_lbl, qn("c:tx"))
            rich = etree.SubElement(tx, qn("c:rich"))

            # body properties
            etree.SubElement(rich, qn("a:bodyPr"))
            etree.SubElement(rich, qn("a:lstStyle"))

            # paragraph
            p = etree.SubElement(rich, qn("a:p"))
            r = etree.SubElement(p, qn("a:r"))

            # 設定字體大小（10pt 清楚可見）
            r_pr = etree.SubElement(r, qn("a:rPr"))
            r_pr.set("lang", "zh-TW")
            r_pr.set("sz", "1000")  # 10pt

            t = etree.SubElement(r, qn("a:t"))
            t.text = name

            # 顯示設定：只顯示自訂文字，不顯示值
            show_val = etree.SubElement(d_lbl, qn("c:showVal"))
            show_val.set("val", "0")
            show_cat = etree.SubElement(d_lbl, qn("c:showCatName"))
            show_cat.set("val", "0")
            show_ser = etree.SubElement(d_lbl, qn("c:showSerName"))
            show_ser.set("val", "0")

    def create_stacked_bar_chart(
        self,
        slide: Slide,
        categories: list[str],
        series_data: dict[str, list[float]],
        title: str = "",
        position: dict = None,
        y_axis_label: str = "",
        y_axis_unit: str = "",
    ):
        """
        建立堆疊直條圖。

        Args:
            slide: 目標投影片
            categories: X 軸類別
            series_data: {系列名稱: [值...]}
            title: 圖表標題
            position: 位置 dict (left, top, width, height)
            y_axis_label: Y 軸標籤
            y_axis_unit: Y 軸單位
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = CategoryChartData()
        chart_data.categories = categories

        for name, values in series_data.items():
            chart_data.add_series(name, values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        for i, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self.colors[i % len(self.colors)]

        # Y 軸設定
        self._set_value_axis_title(chart.value_axis, y_axis_label, y_axis_unit)

        return chart_frame

    def create_pie_chart(
        self,
        slide: Slide,
        categories: list[str],
        values: list[float],
        title: str = "",
        position: dict = None,
    ):
        """
        建立圓餅圖。
        圓餅圖無座標軸，不設定軸標籤。
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("市占率", values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        return chart_frame

    def create_horizontal_bar_chart(
        self,
        slide: Slide,
        categories: list[str],
        series_data: dict[str, list[float]],
        title: str = "",
        position: dict = None,
        y_axis_label: str = "",
        y_axis_unit: str = "",
        highlight_first: bool = True,
        highlight_institution: str = "",
    ):
        """
        建立橫條圖（水平長條圖）。
        適合用於排名展示，類別名稱在左側，數值向右延伸。

        Args:
            slide: 目標投影片
            categories: 類別（通常是機構名稱）
            series_data: {系列名稱: [值...]}
            title: 圖表標題
            position: 位置 dict
            highlight_first: 排名第一是否標金色
            highlight_institution: 目標機構名稱
        """
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        chart_data = CategoryChartData()
        chart_data.categories = categories

        for name, values in series_data.items():
            chart_data.add_series(name, values)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,  # 橫條圖
            position["left"],
            position["top"],
            position["width"],
            position["height"],
            chart_data,
        )

        chart = chart_frame.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # 條件式格式 — 排名第一金色、目標機構深藍
        for i, series in enumerate(chart.series):
            base_color = self.colors[i % len(self.colors)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = base_color

            if len(series_data) == 1:
                for pt_idx, cat in enumerate(categories):
                    point = series.points[pt_idx]
                    if highlight_first and pt_idx == 0:
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0x00)
                    elif highlight_institution and highlight_institution in cat:
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

        # X 軸（數值軸）設定
        self._set_value_axis_title(chart.value_axis, y_axis_label, y_axis_unit)

        return chart_frame

    def create_heatmap(
        self,
        slide: Slide,
        row_labels: list[str],
        col_labels: list[str],
        values: list[list[float]],
        title: str = "",
        position: dict = None,
    ):
        """
        建立熱力圖（以表格 + 條件色彩實現）。
        PowerPoint 不支援原生熱力圖，以彩色表格模擬。
        每個儲存格的背景色根據數值大小從綠（高）到紅（低）漸變。

        Args:
            slide: 目標投影片
            row_labels: 列標籤（如銀行名稱）
            col_labels: 欄標籤（如月份）
            values: 二維數值陣列 [row][col]
            title: 圖表標題
            position: 位置 dict
        """
        from pptx.oxml.ns import qn as _qn

        if position is None:
            position = {"left": Cm(1.5), "top": Cm(4.5), "width": Cm(30.0), "height": Cm(11.0)}

        # 加標題
        if title:
            txBox = slide.shapes.add_textbox(position["left"], Cm(3.8), position["width"], Cm(0.7))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(11)
            p.font.bold = True

        n_rows = len(row_labels) + 1  # +1 for header
        n_cols = len(col_labels) + 1  # +1 for row label column

        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            position["left"], position["top"],
            position["width"], position["height"],
        )
        table = tbl_shape.table

        # 計算數值範圍（用於色彩映射）
        all_vals = [v for row in values for v in row if v is not None]
        min_val = min(all_vals) if all_vals else 0
        max_val = max(all_vals) if all_vals else 1
        val_range = max_val - min_val if max_val != min_val else 1

        # Header row（欄標籤）
        table.cell(0, 0).text = ""
        for col_idx, col_label in enumerate(col_labels):
            table.cell(0, col_idx + 1).text = str(col_label)
            for para in table.cell(0, col_idx + 1).text_frame.paragraphs:
                para.font.size = Pt(8)
                para.font.bold = True

        # Data rows
        for row_idx, row_label in enumerate(row_labels):
            table.cell(row_idx + 1, 0).text = str(row_label)
            for para in table.cell(row_idx + 1, 0).text_frame.paragraphs:
                para.font.size = Pt(8)

            for col_idx, val in enumerate(values[row_idx] if row_idx < len(values) else []):
                cell = table.cell(row_idx + 1, col_idx + 1)
                cell.text = f"{val:.1f}" if val is not None else ""
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(7)

                # 色彩映射：低→紅、中→黃、高→綠
                if val is not None:
                    ratio = (val - min_val) / val_range
                    color = self._heatmap_color(ratio)
                    self._set_table_cell_fill(cell, color)

        return tbl_shape

    def create_quadrant_scatter(
        self,
        slide: Slide,
        data_points: list[dict],
        title: str = "",
        x_label: str = "",
        y_label: str = "",
        position: dict = None,
        x_threshold: float = None,
        y_threshold: float = None,
    ):
        """
        建立風險象限圖（四象限散佈圖）。
        在散佈圖基礎上加入水平線和垂直線標示閾值，將圖分為四個象限。

        Args:
            data_points: [{"name": str, "x": float, "y": float}, ...]
            x_threshold: X 軸分界線（預設為平均值）
            y_threshold: Y 軸分界線（預設為平均值）
        """
        # 先建立基本散佈圖
        chart_frame = self.create_scatter_chart(
            slide, data_points,
            title=title,
            x_label=x_label,
            y_label=y_label,
            position=position,
        )

        # 計算閾值（預設用平均值）
        x_values = [p["x"] for p in data_points if "x" in p]
        y_values = [p["y"] for p in data_points if "y" in p]

        if x_threshold is None and x_values:
            x_threshold = sum(x_values) / len(x_values)
        if y_threshold is None and y_values:
            y_threshold = sum(y_values) / len(y_values)

        # 在圖表區域加入象限分界線說明（用 textbox 模擬）
        if position is None:
            position = {"left": Cm(2), "top": Cm(4), "width": Cm(22), "height": Cm(12)}

        # 四象限標籤（左上、右上、左下、右下）
        quadrant_labels = [
            ("⚠ 高風險\n低規模", Cm(2.5), Cm(4.5)),   # 左上
            ("★ 領先者\n高規模高成長", Cm(18.0), Cm(4.5)),  # 右上
            ("觀察區\n低規模低成長", Cm(2.5), Cm(11.0)),   # 左下
            ("穩定區\n高規模低成長", Cm(18.0), Cm(11.0)),  # 右下
        ]

        for label_text, lx, ly in quadrant_labels:
            txBox = slide.shapes.add_textbox(lx, ly, Cm(5.0), Cm(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label_text
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        return chart_frame

    def _heatmap_color(self, ratio: float) -> RGBColor:
        """
        根據 0~1 的比例回傳熱力圖顏色。
        0.0 = 紅色（低）, 0.5 = 黃色（中）, 1.0 = 綠色（高）。
        """
        ratio = max(0.0, min(1.0, ratio))
        if ratio < 0.5:
            # 紅 → 黃
            r = 0xCC
            g = int(0x33 + (0xCC - 0x33) * (ratio / 0.5))
            b = 0x33
        else:
            # 黃 → 綠
            r = int(0xCC - (0xCC - 0x33) * ((ratio - 0.5) / 0.5))
            g = 0x99
            b = 0x33
        return RGBColor(r, g, b)

    def _set_table_cell_fill(self, cell, color: RGBColor):
        """設定表格儲存格背景色（用於熱力圖）。"""
        from pptx.oxml.ns import qn as _qn

        tc_pr = cell._tc.get_or_add_tcPr()
        for old_fill in tc_pr.findall(_qn("a:solidFill")):
            tc_pr.remove(old_fill)
        solid_fill = etree.SubElement(tc_pr, _qn("a:solidFill"))
        srgb_clr = etree.SubElement(solid_fill, _qn("a:srgbClr"))
        srgb_clr.set("val", str(color))

    # ========== 共用工具方法 ==========

    def _set_value_axis_title(self, value_axis, label: str, unit: str = ""):
        """
        統一設定 Y 軸（數值軸）標題。
        格式：「標籤 (單位)」或僅「標籤」。
        座標軸單位須正確，不混用不同量級。
        """
        if not label:
            return
        value_axis.has_title = True
        axis_text = f"{label} ({unit})" if unit else label
        value_axis.axis_title.text_frame.text = axis_text
