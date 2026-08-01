"""
PowerPoint 簡報生成器
根據 slide_spec 和台新新光金控模板（附件一）生成 16 頁簡報。
所有圖表、表格、文字均為可編輯原生物件，不以圖片嵌入。

模板 Layout 對應：
- 封面/章節分隔: 2_標題投影片 (title only, 居中大字)
- 內容頁(含圖表): 1_標題及內容 (title + body area + page number)
- 章節標題:       2_章節標題 (title 置於下方)
- 結束頁:         3_標題投影片 (title only)
"""

import json
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.slide import Slide

from .template_parser import (
    TemplateParser, TemplateStyle,
    LAYOUT_COVER, LAYOUT_CONTENT, LAYOUT_CHAPTER, LAYOUT_THANKYOU,
)
from .chart_factory import ChartFactory
from .components import (
    fill_executive_summary,
    fill_toc,
    fill_strategy,
    fill_chart,
    fill_table,
    fill_insights_only,
)


class PPTGenerator:
    """根據 slide_spec 生成完整 PowerPoint 簡報。"""

    def __init__(self, template_path: Optional[str] = None):
        self.template_parser = TemplateParser(template_path)
        self.chart_factory = ChartFactory()
        self.style = self.template_parser.get_style()
        self.template_path = template_path

    def generate(self, slide_specs: list[dict], output_path: str) -> str:
        """
        生成完整 16 頁簡報。

        Args:
            slide_specs: 16 頁 slide_spec 陣列
            output_path: 輸出路徑

        Returns:
            輸出檔案路徑
        """
        prs = self.template_parser.create_presentation()

        for spec in slide_specs:
            layout_type = self._resolve_layout_type(spec)
            layout_name = self.template_parser.get_layout_for_page_type(layout_type)
            slide = self._add_slide(prs, layout_name)

            if layout_type == "cover":
                self._build_cover(slide, spec)
            elif layout_type == "chapter":
                self._build_chapter(slide, spec)
            elif layout_type == "thank_you":
                self._build_thank_you(slide, spec)
            elif layout_type == "content":
                self._build_content(slide, spec)

        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        prs.save(output_path)
        return output_path

    def _resolve_layout_type(self, spec: dict) -> str:
        """根據 slide_spec 的 layout 欄位判斷模板 layout 類型。"""
        layout = spec.get("layout", "")
        if layout == "cover":
            return "cover"
        elif layout == "toc":
            return "content"
        elif layout == "chapter_divider":
            return "chapter"
        elif layout == "thank_you":
            return "thank_you"
        elif layout == "executive_summary":
            return "content"
        elif layout == "strategy":
            return "content"
        else:
            # trend_chart, ranking_chart, scatter_chart, comparison_chart, etc.
            return "content"

    def _add_slide(self, prs: Presentation, layout_name: str) -> Slide:
        """新增投影片，使用指定的 layout。"""
        # 尋找對應的 layout
        target_layout = None
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                target_layout = layout
                break

        if target_layout is None:
            # 回退到第一個 layout
            target_layout = prs.slide_layouts[0]

        return prs.slides.add_slide(target_layout)

    def _remove_placeholder(self, slide: Slide, idx: int):
        """從投影片中移除指定的 placeholder（避免顯示空白框）。"""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                sp = ph._element
                sp.getparent().remove(sp)
                return

    def _set_title_style(self, text_frame, font_size=Pt(22), bold=True,
                         color=RGBColor(0x00, 0x33, 0x66), alignment=PP_ALIGN.LEFT):
        """
        統一標題樣式設定。
        所有頁面的標題都透過此方法套用格式，確保一致性。
        修改此處即可全域調整標題外觀，不需逐頁改動。
        """
        for para in text_frame.paragraphs:
            para.font.size = font_size
            para.font.bold = bold
            para.font.color.rgb = color
            para.alignment = alignment

    def _add_page_number(self, slide: Slide, slide_no: int):
        """
        統一頁碼處理：在投影片右下角加入頁碼 textbox。
        封面與結束頁不呼叫此方法，其餘頁面統一使用。
        """
        if not slide_no:
            return
        txBox = slide.shapes.add_textbox(Cm(31.7), Cm(17.8), Cm(1.8), Cm(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_no)
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.RIGHT

    def _add_slide_notes(self, slide: Slide, spec: dict):
        """
        自動生成頁面備註 Notes。
        將 source_ids（資料來源追溯）和補充數據寫入 slide notes，
        簡報者看得到、觀眾看不到，方便 Q&A 時查閱。
        """
        notes_parts = []

        # 資料來源追溯
        source_ids = spec.get("source_ids", [])
        if source_ids:
            notes_parts.append("【資料來源】")
            for sid in source_ids:
                notes_parts.append(f"  • {sid}")

        # 補充 KPI 原始數據
        kpis = spec.get("kpis", [])
        if kpis:
            notes_parts.append("\n【KPI 詳細數據】")
            for kpi in kpis:
                metric_id = kpi.get("metric_id", "")
                label = kpi.get("label", "")
                value = kpi.get("value", "")
                notes_parts.append(f"  • {label}: {value} (metric_id: {metric_id})")

        # 圖表資料摘要
        chart = spec.get("chart", {})
        if chart:
            chart_title = chart.get("title", "")
            notes_parts.append(f"\n【圖表】{chart_title}")

        # 表格資料摘要
        table = spec.get("table", {})
        if table:
            headers = table.get("headers", [])
            row_count = len(table.get("rows", []))
            notes_parts.append(f"\n【表格】{len(headers)} 欄 × {row_count} 行")

        # 寫入 notes
        if notes_parts:
            notes_text = "\n".join(notes_parts)
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    # ========== 封面頁 ==========
    def _build_cover(self, slide: Slide, spec: dict):
        """封面頁 — 使用 2_標題投影片 layout (title placeholder 居中)。"""
        # 移除空的 title placeholder，改用 textbox
        self._remove_placeholder(slide, 0)

        # 主標題
        txBox = slide.shapes.add_textbox(Cm(3.1), Cm(4.0), Cm(27.7), Cm(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = spec.get("title", "信用卡市場分析與經營洞察")
        self._set_title_style(tf, font_size=Pt(32), alignment=PP_ALIGN.CENTER)

        # 副標題
        txBox2 = slide.shapes.add_textbox(Cm(3.1), Cm(7.0), Cm(27.7), Cm(1.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "114 年 1-12 月信用卡業務數據深度解析"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p2.alignment = PP_ALIGN.CENTER

        # 機構名稱
        txBox3 = slide.shapes.add_textbox(Cm(3.1), Cm(9.5), Cm(27.7), Cm(1.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "台新國際商業銀行"
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
        p3.alignment = PP_ALIGN.CENTER

    # ========== 章節分隔頁 ==========
    def _build_chapter(self, slide: Slide, spec: dict):
        """章節分隔頁 — 使用 2_章節標題 layout (title 在下方)。"""
        title = spec.get("title", "")
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text_frame.paragraphs[0].text = title
                self._set_title_style(ph.text_frame, font_size=Pt(28))

        # 統一頁碼
        self._add_page_number(slide, spec.get("slide_no", 0))

    # ========== 結束頁 ==========
    def _build_thank_you(self, slide: Slide, spec: dict):
        """結束頁 — 使用 3_標題投影片 layout。"""
        # 移除空的 title placeholder
        self._remove_placeholder(slide, 0)

        txBox = slide.shapes.add_textbox(Cm(3.1), Cm(5.0), Cm(27.7), Cm(2.5))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "感謝聆聽"
        self._set_title_style(tf, font_size=Pt(36), alignment=PP_ALIGN.CENTER)

        txBox2 = slide.shapes.add_textbox(Cm(3.1), Cm(8.5), Cm(27.7), Cm(2.0))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "台新國際商業銀行 — 信用卡市場分析與經營洞察"
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf2.add_paragraph()
        p3.text = "資料來源：各金融機構信用卡重要資訊揭露（114年1-12月）"
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

        # Excel 原始資料連結（Data Lineage 追溯入口）
        excel_path = spec.get("excel_source", "") or self._get_excel_source()
        if excel_path:
            txBox3 = slide.shapes.add_textbox(Cm(3.1), Cm(12.0), Cm(27.7), Cm(1.0))
            tf3 = txBox3.text_frame
            p4 = tf3.paragraphs[0]
            p4.alignment = PP_ALIGN.CENTER
            run = p4.add_run()
            run.text = f"📎 原始資料：{excel_path}"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
            run.font.underline = True
            # 加入超連結
            run.hyperlink.address = excel_path

    def _get_excel_source(self) -> str:
        """取得 Excel 來源檔案路徑（從模板路徑推斷或使用預設）。"""
        import os
        # 預設指向附件四
        default = "附件四_預期修正參照資料.xlsx"
        if os.path.exists(default):
            return default
        return ""

    # ========== 內容頁 ==========
    def _build_content(self, slide: Slide, spec: dict):
        """
        內容頁 — 使用 1_標題及內容 layout。
        Placeholders:
          idx=0 TITLE: left=1.0 top=0.7 w=31.3 h=2.1
          idx=1 BODY:  left=1.0 top=3.5 w=31.2 h=13.6
          idx=12 SLIDE_NUMBER: left=31.7 top=17.8
        """
        title = spec.get("title", "")
        headline = spec.get("headline", "")
        slide_no = spec.get("slide_no", 0)

        # 填入 TITLE placeholder，移除 SLIDE_NUMBER placeholder（改用統一 textbox）
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = title
                self._set_title_style(ph.text_frame)

        # 移除原生 SLIDE_NUMBER placeholder，改用統一的 textbox 頁碼
        self._remove_placeholder(slide, 12)

        # 移除 BODY placeholder（我們用自訂 shapes 取代）
        self._remove_placeholder(slide, 1)

        # 統一頁碼
        self._add_page_number(slide, slide_no)

        # 自動生成頁面備註 Notes（資料來源追溯 + 補充資訊）
        self._add_slide_notes(slide, spec)

        # Headline (核心訊息)
        if headline:
            txBox = slide.shapes.add_textbox(Cm(1.0), Cm(2.9), Cm(31.2), Cm(1.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"◆ {headline}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # 根據頁面內容類型決定如何填充
        layout = spec.get("layout", "")

        if layout == "executive_summary":
            self._fill_executive_summary(slide, spec)
        elif layout == "toc":
            self._fill_toc(slide, spec)
        elif layout == "strategy":
            self._fill_strategy(slide, spec)
        elif spec.get("chart"):
            self._fill_chart(slide, spec)
        elif spec.get("table"):
            self._fill_table(slide, spec)
        else:
            self._fill_insights_only(slide, spec)

    def _fill_executive_summary(self, slide: Slide, spec: dict):
        """填充 Executive Summary 頁 — 委派給 components 模組。"""
        fill_executive_summary(slide, spec)

    def _fill_toc(self, slide: Slide, spec: dict):
        """填充目錄頁 — 委派給 components 模組。"""
        fill_toc(slide, spec)

    def _fill_strategy(self, slide: Slide, spec: dict):
        """填充策略建議頁 — 委派給 components 模組。"""
        fill_strategy(slide, spec)

    def _fill_chart(self, slide: Slide, spec: dict):
        """填充含圖表的頁面 — 委派給 components 模組。"""
        fill_chart(slide, spec, self.chart_factory)

    def _fill_table(self, slide: Slide, spec: dict):
        """填充含原生表格的頁面 — 委派給 components 模組。"""
        fill_table(slide, spec)

    def _fill_insights_only(self, slide: Slide, spec: dict):
        """只有洞察文字的頁面 — 委派給 components 模組。"""
        fill_insights_only(slide, spec)
